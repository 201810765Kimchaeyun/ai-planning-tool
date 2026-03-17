import streamlit as st
import streamlit.components.v1 as components
from pptx import Presentation
from bs4 import BeautifulSoup
import anthropic
import json
import re
import pandas as pd
import io
# ─────────────────────────────────────────────
# Page config
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="AI 기획 리뷰 툴",
    page_icon="🤖",
    layout="wide",
)
# ─────────────────────────────────────────────
# Leave warning
# ─────────────────────────────────────────────
def inject_leave_warning(enabled: bool):
    script = f"""
    <script>
    const warningEnabled = {str(enabled).lower()};
    window.onbeforeunload = null;
    if (warningEnabled) {{
        window.onbeforeunload = function (e) {{
            e.preventDefault();
            e.returnValue = '';
            return '';
        }};
    }}
    </script>
    """
    components.html(script, height=0)
# ─────────────────────────────────────────────
# 화면정의서 HTML 파서 (scenes JSON 방식)
# ─────────────────────────────────────────────
def parse_screen_definition_html(content: str) -> list[dict]:
    """화면정의서 HTML에서 var D = {...} JSON을 추출해 각 씬의 스펙을 파싱"""
    pages = []
    # var D = {...} 패턴 추출
    match = re.search(r'var\s+D\s*=\s*(\{.*?\});\s*\n', content, re.DOTALL)
    if not match:
        # fallback: 단순 BeautifulSoup 텍스트 추출
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        pages.append({"title": "전체 내용", "text": text[:8000]})
        return pages
    try:
        data = json.loads(match.group(1))
        scenes = data.get("scenes", {})
        for scene_id, scene in scenes.items():
            title = scene.get("title", scene_id)
            table_html = scene.get("tableHtml", "")
            memo = scene.get("memo", "")
            # tableHtml(스펙 정의)에서 텍스트 추출
            spec_lines = []
            if table_html:
                soup = BeautifulSoup(table_html, "html.parser")
                # spec-title-row에서 요소명 추출
                for title_row in soup.find_all("tr", class_="spec-title-row"):
                    td = title_row.find("td")
                    if td:
                        badge = td.find("span", class_="badge-inline")
                        if badge:
                            badge.decompose()
                        spec_lines.append(f"[요소] {td.get_text(strip=True)}")
                # 일반 행 텍스트 추출
                for row in soup.find_all("tr"):
                    if "spec-title-row" in row.get("class", []):
                        continue
                    th = row.find("th")
                    tds = row.find_all("td")
                    if th and tds:
                        th_text = th.get_text(strip=True).replace("⠿", "").strip()
                        td_text = " ".join(td.get_text(strip=True) for td in tds)
                        if th_text and td_text:
                            spec_lines.append(f"  {th_text}: {td_text}")
            if memo:
                spec_lines.append(f"  [메모] {memo}")
            if spec_lines:
                pages.append({
                    "title": title,
                    "text": "\n".join(spec_lines)
                })
    except Exception as e:
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        pages.append({"title": "전체 내용 (파싱 오류 fallback)", "text": text[:8000]})
    return pages
# ─────────────────────────────────────────────
# 일반 HTML 파서 (정책 문서용)
# ─────────────────────────────────────────────
def parse_general_html(content: str) -> str:
    """일반 HTML에서 텍스트 추출 (정책 문서 등)"""
    soup = BeautifulSoup(content, "html.parser")
    # script, style 제거
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)
# ─────────────────────────────────────────────
# PPTX 파서
# ─────────────────────────────────────────────
def parse_pptx(file) -> list[dict]:
    prs = Presentation(file)
    pages = []
    for i, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if shape.shape_type == 19:  # TABLE
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        text += " | ".join(row_texts) + "\n"
            elif hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
        if text.strip():
            pages.append({"title": f"Slide {i+1}", "text": text.strip()})
    return pages
# ─────────────────────────────────────────────
# Notion URL 텍스트 추출
# ─────────────────────────────────────────────
def fetch_notion_text(url: str) -> str:
    """노션 페이지에서 텍스트 추출 (공개 페이지 한정)"""
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        resp = requests.get(url, headers=headers, timeout=10)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)[:10000]
    except Exception as e:
        return f"[노션 페이지 접근 오류: {e}]"
# ─────────────────────────────────────────────
# XLSX 내보내기 (시트별 분리)
# ─────────────────────────────────────────────
def build_xlsx_report(results: dict) -> bytes:
    """결과를 시트별로 분리한 xlsx 파일 생성"""
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        sheet_map = {
            "ux_issues": {
                "sheet": "UX Issues",
                "cols": ["Title", "Description", "Impact"],
                "keys": ["title", "description", "impact"],
            },
            "policy_gaps": {
                "sheet": "Policy Gaps",
                "cols": ["Title", "Description", "Recommendation"],
                "keys": ["title", "description", "recommendation"],
            },
            "edge_cases": {
                "sheet": "Edge Cases",
                "cols": ["Title", "Scenario", "Expected Behavior"],
                "keys": ["title", "scenario", "expected_behavior"],
            },
            "test_cases": {
                "sheet": "QA Test Cases",
                "cols": ["TC #", "Screen", "Test Case", "Description", "Expected Result", "Priority"],
                "keys": ["id", "screen", "test_case", "description", "expected_result", "priority"],
            },
        }
        for key, meta in sheet_map.items():
            items = results.get(key, [])
            if not items:
                continue
            df = pd.DataFrame([
                {c: item.get(k, "") for c, k in zip(meta["cols"], meta["keys"])}
                for item in items
            ])
            df.to_excel(writer, sheet_name=meta["sheet"], index=False)
            # 열 너비 자동 조정
            ws = writer.sheets[meta["sheet"]]
            for col in ws.columns:
                max_len = max((len(str(cell.value or "")) for cell in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 60)
    return output.getvalue()
# ─────────────────────────────────────────────
# HTML 리포트 빌더
# ─────────────────────────────────────────────
def build_html_report(results: dict) -> str:
    sections_html = ""
    label_map = {
        "ux_issues": ("1. UX Issues", ["Title", "Description", "Impact"]),
        "policy_gaps": ("2. Policy Gaps", ["Title", "Description", "Recommendation"]),
        "edge_cases": ("3. Edge Cases", ["Title", "Scenario", "Expected Behavior"]),
        "test_cases": ("4. QA Test Cases", ["TC #", "Test Case", "Description", "Expected Result", "Priority"]),
    }
    key_map = {
        "ux_issues": ["title", "description", "impact"],
        "policy_gaps": ["title", "description", "recommendation"],
        "edge_cases": ["title", "scenario", "expected_behavior"],
        "test_cases": ["id", "test_case", "description", "expected_result", "priority"],
    }
    for key, (heading, cols) in label_map.items():
        items = results.get(key, [])
        if not items:
            continue
        rows = ""
        for item in items:
            cells = "".join(f"<td>{item.get(k, '')}</td>" for k in key_map[key])
            rows += f"<tr>{cells}</tr>"
        header = "".join(f"<th>{c}</th>" for c in cols)
        sections_html += f"""
        <h2>{heading}</h2>
        <table>
            <thead><tr>{header}</tr></thead>
            <tbody>{rows}</tbody>
        </table>
        """
    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <title>AI 기획 리뷰 결과</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 32px; color: #1e293b; line-height: 1.6; }}
        h1 {{ margin-bottom: 8px; }}
        h2 {{ margin-top: 32px; border-bottom: 2px solid #e2e8f0; padding-bottom: 6px; }}
        table {{ width: 100%; border-collapse: collapse; margin-top: 12px; margin-bottom: 28px; table-layout: fixed; word-wrap: break-word; }}
        th, td {{ border: 1px solid #cbd5e1; padding: 10px; text-align: left; vertical-align: top; font-size: 14px; }}
        th {{ background: #f1f5f9; }}
        .meta {{ color: #64748b; margin-bottom: 24px; }}
    </style>
</head>
<body>
    <h1>AI 기획 리뷰 결과</h1>
    <div class="meta">Generated by AI Planning Review Tool</div>
    {sections_html}
</body>
</html>"""
# ─────────────────────────────────────────────
# Custom CSS
# ─────────────────────────────────────────────
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Sans+KR:wght@400;500;600&display=swap');
    html, body, [class*="css"] { font-family: 'IBM Plex Sans KR', sans-serif; }
    .main-header {
        background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
        color: #f8fafc; padding: 2rem 2.5rem; border-radius: 12px;
        margin-bottom: 2rem; border-left: 4px solid #6366f1;
    }
    .main-header h1 { margin: 0; font-size: 1.8rem; font-weight: 600; }
    .main-header p  { margin: 0.4rem 0 0; color: #94a3b8; font-size: 0.9rem; }
    .section-title {
        font-size: 1.05rem; font-weight: 600; color: #1e293b;
        border-bottom: 2px solid #e2e8f0; padding-bottom: 0.5rem;
        margin-bottom: 1rem; margin-top: 1.5rem;
    }
    .section-box {
        background: #f8fafc; border: 1px solid #e2e8f0;
        border-radius: 10px; padding: 1.2rem 1.5rem; margin-bottom: 1rem;
    }
    .section-box h4 { margin: 0 0 0.5rem; font-size: 0.9rem; color: #6366f1; font-weight: 600; }
    .slide-card {
        background: #f8fafc; border: 1px solid #e2e8f0;
        border-radius: 8px; padding: 1rem 1.2rem; margin-bottom: 0.8rem;
        font-size: 0.85rem; line-height: 1.7;
    }
    .stat-box {
        background: #f1f5f9; border-radius: 8px; padding: 1rem; text-align: center;
    }
    .stat-box .num { font-size: 2rem; font-weight: 700; color: #6366f1; }
    .stat-box .lbl { font-size: 0.8rem; color: #64748b; margin-top: 2px; }
    .policy-box {
        background: #fffbeb; border: 1px solid #fde68a;
        border-radius: 10px; padding: 1rem 1.2rem; margin-bottom: 0.5rem;
    }
    .policy-box h4 { margin: 0 0 0.3rem; font-size: 0.85rem; color: #92400e; }
    .block-container {
        max-width: 1200px !important;
    }
</style>
""", unsafe_allow_html=True)
# ─────────────────────────────────────────────
# Session state 초기화
# ─────────────────────────────────────────────
if "current_results" not in st.session_state:
    st.session_state.current_results = {}
if "current_file_name" not in st.session_state:
    st.session_state.current_file_name = ""
if "current_tc_strategy" not in st.session_state:
    st.session_state.current_tc_strategy = "regression"
# ─────────────────────────────────────────────
# Sidebar
# ─────────────────────────────────────────────
with st.sidebar:
    st.markdown("### 🔑 API 설정")
    api_key = st.text_input(
        "Anthropic API Key", type="password",
        placeholder="sk-ant-api03-...",
        help="https://console.anthropic.com 에서 발급",
    )
    if api_key:
        st.success("API 키 입력 완료", icon="✅")
    else:
        st.warning("API 키를 입력해야 리뷰 결과 생성이 가능합니다.", icon="⚠️")
    st.markdown("---")
    st.markdown("**API 키 발급 방법**")
    st.markdown("1. [console.anthropic.com](https://console.anthropic.com) 접속")
    st.markdown("2. **API Keys** 메뉴 → **Create Key**")
    st.markdown("3. 발급된 키를 위 입력란에 붙여넣기")
# ─────────────────────────────────────────────
# Header
# ─────────────────────────────────────────────
st.markdown("""
<div class="main-header">
    <h1>🤖 AI 기획 리뷰 툴</h1>
    <p>화면정의서(.html / .pptx)를 업로드하면 AI가 UX Issues, Policy Gaps, Edge Cases, QA Test Cases를 자동 생성합니다.</p>
</div>
""", unsafe_allow_html=True)
# ─────────────────────────────────────────────
# STEP 1: 기획안 파일 업로드
# ─────────────────────────────────────────────
st.markdown('<div class="section-title">📁 STEP 1 — 기획안 파일 업로드</div>', unsafe_allow_html=True)
st.caption("화면 단위 기획안 파일을 올려주세요. HTML(화면정의서) 또는 PPT 형식 지원.")
uploaded_file = st.file_uploader(
    "기획안 파일 업로드 (.html 또는 .pptx)",
    type=["html", "htm", "pptx"],
    key="main_file",
)
inject_leave_warning(uploaded_file is not None)
main_pages = []
all_text = ""
if uploaded_file:
    file_name = uploaded_file.name.lower()
    raw_content = uploaded_file.read()
    if file_name.endswith(".pptx"):
        import io
        main_pages = parse_pptx(io.BytesIO(raw_content))
    else:
        html_content = raw_content.decode("utf-8", errors="ignore")
        main_pages = parse_screen_definition_html(html_content)
    all_text = "\n\n".join(
        f"=== {p['title']} ===\n{p['text']}" for p in main_pages
    )
    with st.expander(f"📄 추출된 내용 보기 ({len(main_pages)}개 화면)", expanded=False):
        for p in main_pages:
            st.markdown(f"**{p['title']}**")
            st.markdown(
                '<div class="slide-card">' + p["text"].replace("\n", "<br>") + '</div>',
                unsafe_allow_html=True,
            )
# ─────────────────────────────────────────────
# STEP 2: 공통 정책 문서 (선택)
# ─────────────────────────────────────────────
st.markdown('<div class="section-title">📋 STEP 2 — 공통 정책 문서 (선택)</div>', unsafe_allow_html=True)
st.caption("화면 단위가 아닌 서비스 공통 정책 문서를 추가하면 더 정확한 분석이 가능합니다.")
policy_text = ""
col_policy_file, col_policy_paste = st.columns([1, 1])
with col_policy_file:
    st.markdown('<div class="section-box"><h4>📎 파일 첨부</h4>', unsafe_allow_html=True)
    policy_file = st.file_uploader(
        "정책 문서 파일 (.html, .htm, .docx)",
        type=["html", "htm", "docx"],
        key="policy_file",
        label_visibility="collapsed",
    )
    st.markdown('</div>', unsafe_allow_html=True)
with col_policy_paste:
    st.markdown('<div class="section-box"><h4>📝 직접 붙여넣기</h4>', unsafe_allow_html=True)
    st.caption("노션, 컨플루언스 등에서 텍스트를 복사해 붙여넣으세요.")
    policy_paste = st.text_area(
        "정책 내용 직접 입력",
        placeholder="노션 페이지 내용, 정책 문서 텍스트 등을 여기에 붙여넣으세요...",
        height=160,
        key="policy_paste",
        label_visibility="collapsed",
    )
    st.markdown('</div>', unsafe_allow_html=True)
# 정책 파일 처리
if policy_file:
    pf_name = policy_file.name.lower()
    pf_content = policy_file.read()
    if pf_name.endswith(".docx"):
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(pf_content))
            policy_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            st.warning(f"DOCX 파싱 오류: {e}")
    else:
        html_str = pf_content.decode("utf-8", errors="ignore")
        policy_text = parse_general_html(html_str)
    if policy_text:
        st.success(f"정책 문서 로드 완료 ({len(policy_text):,}자)", icon="✅")
        with st.expander("정책 문서 내용 미리보기", expanded=False):
            st.text(policy_text[:3000] + ("..." if len(policy_text) > 3000 else ""))
# 직접 붙여넣기 처리
if policy_paste and policy_paste.strip():
    pasted = policy_paste.strip()
    policy_text = (policy_text + "\n\n" + pasted).strip() if policy_text else pasted
    st.success(f"텍스트 입력 완료 ({len(pasted):,}자)", icon="✅")
# ─────────────────────────────────────────────
# STEP 3: 분석 항목 선택
# ─────────────────────────────────────────────
st.markdown('<div class="section-title">⚙️ STEP 3 — 분석 항목 선택</div>', unsafe_allow_html=True)
st.caption("생성할 항목을 선택하세요. 선택한 항목만 API를 호출합니다.")
col_a, col_b, col_c, col_d = st.columns(4)
with col_a:
    do_ux = st.checkbox("🎨 UX Issues", value=True)
with col_b:
    do_policy = st.checkbox("📜 Policy Gaps", value=True)
with col_c:
    do_edge = st.checkbox("🔍 Edge Cases", value=True)
with col_d:
    do_test = st.checkbox("🧪 QA Test Cases", value=True)
selected = [k for k, v in {
    "ux_issues": do_ux,
    "policy_gaps": do_policy,
    "edge_cases": do_edge,
    "test_cases": do_test,
}.items() if v]
if not selected:
    st.warning("최소 1개 이상의 분석 항목을 선택해주세요.")
# QA Test Cases 선택 시 테스트 전략 선택
test_strategy = "regression"
if do_test:
    st.markdown("")
    st.markdown("**🎯 QA 테스트 전략 선택**")
    st.caption("인력/일정 상황에 맞게 테스트 깊이를 선택하세요.")
    strategy_col1, strategy_col2, strategy_col3 = st.columns(3)
    with strategy_col1:
        st.markdown("""
        <div style="border:2px solid #ef4444;border-radius:8px;padding:0.8rem;min-height:100px;">
            <div style="font-weight:700;color:#ef4444;">🔴 Smoke Test</div>
            <div style="font-size:0.78rem;color:#64748b;margin-top:4px;">
                핵심 플로우만 빠르게<br>
                · 서비스 중단 케이스<br>
                · 결제·인증·데이터 핵심 경로<br>
                · 긴급 릴리즈·핫픽스용
            </div>
        </div>
        """, unsafe_allow_html=True)
    with strategy_col2:
        st.markdown("""
        <div style="border:2px solid #f59e0b;border-radius:8px;padding:0.8rem;min-height:100px;">
            <div style="font-weight:700;color:#f59e0b;">🟡 Regression</div>
            <div style="font-size:0.78rem;color:#64748b;margin-top:4px;">
                주요 기능 + 경계값<br>
                · 기능별 정상/비정상<br>
                · 주요 엣지케이스<br>
                · 일반 스프린트용 (기본값)
            </div>
        </div>
        """, unsafe_allow_html=True)
    with strategy_col3:
        st.markdown("""
        <div style="border:2px solid #10b981;border-radius:8px;padding:0.8rem;min-height:100px;">
            <div style="font-weight:700;color:#10b981;">🟢 Full Coverage</div>
            <div style="font-size:0.78rem;color:#64748b;margin-top:4px;">
                가능한 모든 케이스<br>
                · 권한·네트워크·오류 처리<br>
                · 마이너 엣지케이스 포함<br>
                · 대규모 릴리즈·감사용
            </div>
        </div>
        """, unsafe_allow_html=True)
    st.markdown("")
    test_strategy = st.radio(
        "전략 선택",
        options=["smoke", "regression", "full"],
        format_func=lambda x: {"smoke": "🔴 Smoke Test", "regression": "🟡 Regression", "full": "🟢 Full Coverage"}[x],
        index=1,
        horizontal=True,
        label_visibility="collapsed",
    )
# ─────────────────────────────────────────────
# STEP 4: AI 리뷰 생성
# ─────────────────────────────────────────────
st.divider()
st.markdown('<div class="section-title">🚀 STEP 4 — AI 리뷰 결과 생성</div>', unsafe_allow_html=True)
can_run = uploaded_file and all_text and selected
if st.button(
    "✨ 리뷰 결과 생성하기",
    type="primary",
    use_container_width=True,
    disabled=not can_run,
):
    if not api_key:
        st.error("사이드바에 Anthropic API 키를 먼저 입력해주세요.")
        st.stop()
    if not can_run:
        st.error("기획안 파일을 업로드하고 분석 항목을 선택해주세요.")
        st.stop()
    # ── 분석 항목별 JSON 키 & 프롬프트 정의 ──
    item_configs = {
        "ux_issues": {
            "label": "UX Issues",
            "schema": '{"ux_issues": [{"title": "...", "description": "...", "impact": "..."}]}',
        },
        "policy_gaps": {
            "label": "Policy Gaps",
            "schema": '{"policy_gaps": [{"title": "...", "description": "...", "recommendation": "..."}]}',
        },
        "edge_cases": {
            "label": "Edge Cases",
            "schema": '{"edge_cases": [{"title": "...", "scenario": "...", "expected_behavior": "..."}]}',
        },
        "test_cases": {
            "label": "QA Test Cases",
            # ▼ [변경 1] screen 필드 추가
            "schema": '{"test_cases": [{"id": "TC-001", "screen": "화면명", "test_case": "...", "description": "...", "expected_result": "...", "priority": "High"}]}',
            "strategy": test_strategy,
        },
    }
    # ── 컨텍스트 구성 ──
    policy_section = ""
    if policy_text:
        policy_section = f"\n\n[공통 정책 문서]\n{policy_text[:5000]}"
    main_content = all_text[:10000]
    # ── API 호출 (선택 항목만) ──
    results = {}
    client = anthropic.Anthropic(api_key=api_key)
    progress = st.progress(0, text="분석 준비 중...")
    total = len(selected)
    for idx, key in enumerate(selected):
        cfg = item_configs[key]
        progress.progress((idx) / total, text=f"📊 {cfg['label']} 분석 중... ({idx+1}/{total})")
        # test_cases는 전략별 프롬프트, 나머지는 공통 프롬프트
        if key == "test_cases":
            strategy_key = cfg.get("strategy", "regression")
            strategy_guide = {
                # ▼ [변경 1] 각 전략에 screen 필드 안내 추가
                "smoke": """priority는 반드시 High / Medium / Low 중 하나만 사용하세요.
screen 필드에는 기획안에서 === 화면명 === 형식으로 표시된 화면명을 정확히 입력하세요.
테스트 전략: Smoke Test — 서비스 중단을 일으킬 수 있는 Critical Path만 뽑아라.
기준:
- 핵심 사용자 플로우 (회원가입/로그인/결제/핵심 기능 진입 등)
- 이게 안 되면 서비스 자체가 동작 불가한 케이스
- 단순하고 명확한 케이스 위주
기능이 단순하면 2~3개, 복잡한 화면도 5개를 넘기지 마라.
우선순위는 대부분 High로 설정하라.""",
                "regression": """priority는 반드시 High / Medium / Low 중 하나만 사용하세요.
screen 필드에는 기획안에서 === 화면명 === 형식으로 표시된 화면명을 정확히 입력하세요.
테스트 전략: Regression Test — 주요 기능의 정상/비정상 케이스와 주요 경계값을 뽑아라.
기준:
- 각 기능의 정상 동작 케이스
- 입력값 경계값 (빈값, 최대값, 특수문자 등)
- 자주 발생하는 비정상 케이스 (권한 오류, 중복 처리 등)
- 마이너한 UI 케이스나 극히 드문 에러 상황은 제외
기능 복잡도에 비례해서 작성하라. 단순한 화면은 3~5개, 복잡한 화면은 10~15개가 적절하다.""",
                "full": """priority는 반드시 High / Medium / Low 중 하나만 사용하세요.
screen 필드에는 기획안에서 === 화면명 === 형식으로 표시된 화면명을 정확히 입력하세요.
테스트 전략: Full Coverage — 가능한 모든 분기와 케이스를 뽑아라.
기준:
- 모든 정상/비정상 케이스
- 경계값, 극단값, 특수문자, 빈값 등 입력 케이스 전체
- 권한별 접근 제어, 네트워크 오류, 타임아웃
- 동시성 이슈, 중복 요청, 세션 만료 등 엣지케이스
- 마이너한 UI 상태, 로딩/빈 화면/에러 화면 처리
각 기능을 철저하게 분석해 누락 없이 작성하라. 기능이 복잡할수록 더 많이 작성하라.""",
            }[strategy_key]
            system_prompt = (
                f"당신은 시니어 QA 엔지니어입니다.\n"
                f"서비스 기획안을 분석하여 QA Test Cases만 도출하세요.\n\n"
                f"반드시 아래 JSON 형식으로만 응답하세요. 다른 텍스트 없이 JSON만 출력하세요.\n\n"
                f"{cfg['schema']}\n\n"
                f"{strategy_guide}"
            )
        else:
            system_prompt = (
                f"당신은 시니어 서비스 기획 리뷰어입니다.\n"
                f"서비스 기획안을 분석하여 **{cfg['label']}** 항목만 도출하세요.\n\n"
                f"반드시 아래 JSON 형식으로만 응답하세요. 다른 텍스트 없이 JSON만 출력하세요.\n\n"
                f"{cfg['schema']}\n\n"
                f"기획안의 복잡도와 내용에 비례해서 필요한 만큼만 작성하세요. 각 항목은 구체적으로 작성하세요."
            )
        user_prompt = (
            f"다음 서비스 기획안을 분석하여 {cfg['label']}를 JSON으로 생성하세요."
            f"{policy_section}\n\n"
            f"[기획안 내용]\n{main_content}"
        )
        try:
            # Full Coverage TC는 화면별 분할 호출 후 합산
            if key == "test_cases" and cfg.get("strategy") == "full" and len(main_pages) > 1:
                all_tcs = []
                screen_total = len(main_pages)
                for s_idx, page in enumerate(main_pages):
                    progress.progress(
                        (idx + (s_idx / screen_total)) / total,
                        text=f"🔬 Full Coverage: {page['title']} 분석 중... ({s_idx+1}/{screen_total})"
                    )
                    screen_prompt = (
                        f"다음 서비스 기획안의 [{page['title']}] 화면을 분석하여 "
                        f"QA Test Cases를 JSON으로 생성하세요.\n\n"
                        f"{policy_section}\n\n"
                        f"[화면 내용]\n{page['text']}"
                    )
                    parse_errors = []
                    try:
                        r = client.messages.create(
                            model="claude-sonnet-4-20250514",
                            max_tokens=6000,
                            system=system_prompt,
                            messages=[{"role": "user", "content": screen_prompt}],
                        )
                        raw = r.content[0].text.strip()
                        raw_original = raw
                        # ── JSON 펜스 제거 ──
                        if raw.startswith("```"):
                            raw = re.sub(r"^```(?:json)?\s*", "", raw)
                            raw = re.sub(r"\s*```$", "", raw)
                        raw = raw.strip()
                        # { 시작 ~ } 끝으로 잘라내기
                        brace_start = raw.find("{")
                        brace_end   = raw.rfind("}")
                        if brace_start != -1 and brace_end != -1:
                            raw = raw[brace_start:brace_end+1]
                        # ── 정상 파싱 시도 ──
                        try:
                            data = json.loads(raw)
                            screen_tcs = data.get("test_cases", [])
                        except json.JSONDecodeError:
                            screen_tcs = []
                            arr_start = raw.find('"test_cases"')
                            if arr_start != -1:
                                bracket = raw.find("[", arr_start)
                                if bracket != -1:
                                    partial = raw[bracket:]
                                    for m in re.finditer(r'\{[^{}]*\}', partial, re.DOTALL):
                                        try:
                                            obj = json.loads(m.group())
                                            if "test_case" in obj or "id" in obj:
                                                screen_tcs.append(obj)
                                        except Exception:
                                            continue
                        prefix = f"TC-{s_idx+1:02d}"
                        for i, tc in enumerate(screen_tcs):
                            tc["id"] = f"{prefix}-{i+1:03d}"
                            tc["screen"] = page["title"]  # 화면명 명시적 덮어쓰기
                        all_tcs.extend(screen_tcs)
                    except json.JSONDecodeError as e:
                        parse_errors.append({
                            "screen": page["title"],
                            "error": str(e),
                            "raw": raw_original[:500],
                        })
                    except Exception as e:
                        parse_errors.append({
                            "screen": page["title"],
                            "error": str(e),
                            "raw": "",
                        })
                results[key] = all_tcs
                if parse_errors:
                    st.session_state["tc_parse_errors"] = parse_errors
                else:
                    st.session_state.pop("tc_parse_errors", None)
            else:
                # 단일 호출 (Smoke / Regression / 화면 1개)
                response = client.messages.create(
                    model="claude-sonnet-4-20250514",
                    max_tokens=8000,
                    system=system_prompt,
                    messages=[{"role": "user", "content": user_prompt}],
                )
                raw = response.content[0].text.strip()
                if "```" in raw:
                    parts = raw.split("```")
                    if len(parts) > 1:
                        raw = parts[1]
                    if raw.startswith("json"):
                        raw = raw[4:]
                data = json.loads(raw.strip())
                results[key] = data.get(key, [])
        except Exception as e:
            st.warning(f"{cfg['label']} 파싱 오류: {e}")
            results[key] = []
    progress.progress(1.0, text="✅ 분석 완료!")
    # 세션에 결과 저장
    st.session_state.current_results = results
    st.session_state.current_file_name = uploaded_file.name if uploaded_file else "unknown"
    st.session_state.current_tc_strategy = test_strategy
    st.rerun()
# ─────────────────────────────────────────────
# 결과 렌더링 함수
# ─────────────────────────────────────────────
SECTION_META = {
    "ux_issues": {
        "label": "🎨 1. UX Issues",
        "cols": ["Title", "Description", "Impact"],
        "keys": ["title", "description", "impact"],
        "col_config": {
            "Title":       st.column_config.TextColumn(width="medium", max_chars=None),
            "Description": st.column_config.TextColumn(width="large",  max_chars=None),
            "Impact":      st.column_config.TextColumn(width="medium", max_chars=None),
        },
    },
    "policy_gaps": {
        "label": "📜 2. Policy Gaps",
        "cols": ["Title", "Description", "Recommendation"],
        "keys": ["title", "description", "recommendation"],
        "col_config": {
            "Title":          st.column_config.TextColumn(width="medium", max_chars=None),
            "Description":    st.column_config.TextColumn(width="large",  max_chars=None),
            "Recommendation": st.column_config.TextColumn(width="large",  max_chars=None),
        },
    },
    "edge_cases": {
        "label": "🔍 3. Edge Cases",
        "cols": ["Title", "Scenario", "Expected Behavior"],
        "keys": ["title", "scenario", "expected_behavior"],
        "col_config": {
            "Title":            st.column_config.TextColumn(width="medium", max_chars=None),
            "Scenario":         st.column_config.TextColumn(width="large",  max_chars=None),
            "Expected Behavior":st.column_config.TextColumn(width="large",  max_chars=None),
        },
    },
}
TC_COLS = ["TC #", "Screen", "Test Case", "Description", "Expected Result", "Priority"]
TC_KEYS = ["id", "screen", "test_case", "description", "expected_result", "priority"]
TC_COL_CONFIG = {
    "TC #":           st.column_config.TextColumn(width="small",  max_chars=None),
    "Screen":         st.column_config.TextColumn(width="small",  max_chars=None),
    "Test Case":      st.column_config.TextColumn(width="medium", max_chars=None),
    "Description":    st.column_config.TextColumn(width="large",  max_chars=None),
    "Expected Result":st.column_config.TextColumn(width="large",  max_chars=None),
    "Priority":       st.column_config.SelectboxColumn(options=["High","Medium","Low"], width="small"),
}
def items_to_df(items: list, cols: list, keys: list) -> pd.DataFrame:
    return pd.DataFrame([{c: item.get(k, "") for c, k in zip(cols, keys)} for item in items])
def df_to_items(df: pd.DataFrame, cols: list, keys: list) -> list:
    return [{k: row.get(c, "") for c, k in zip(cols, keys)} for row in df.to_dict("records")]
def is_dirty(original: list, edited_df: pd.DataFrame, cols: list, keys: list) -> bool:
    edited = df_to_items(edited_df, cols, keys)
    return edited != original

# ▼ [변경 2] 행 수에 맞게 높이를 계산하는 헬퍼 (기존 행 + 빈 행 1개)
def calc_height(n_rows: int, max_height: int = 800) -> int:
    """헤더(40px) + 데이터 행(38px each) + 빈 행 1개(38px)"""
    return min(max(100, 40 + (n_rows + 1) * 38), max_height)

def render_results(results: dict, file_name: str):
    buf_key = "edit_buf"
    if buf_key not in st.session_state:
        st.session_state[buf_key] = json.loads(json.dumps(results))
    buf = st.session_state[buf_key]
    LABEL_DISPLAY = {
        "ux_issues": "🎨 UX Issues",
        "policy_gaps": "📜 Policy Gaps",
        "edge_cases": "🔍 Edge Cases",
        "test_cases": "🧪 QA Test Cases",
    }
    st.markdown(f"""
    <div style="background:#f0f9ff;border:1px solid #bae6fd;border-radius:10px;
                padding:1rem 1.5rem;margin-bottom:1rem;">
        <div style="font-size:0.8rem;color:#0369a1;font-weight:600;">📂 분석 파일</div>
        <div style="font-size:1rem;font-weight:700;color:#0c4a6e;">{file_name}</div>
    </div>
    """, unsafe_allow_html=True)
    active_keys = [k for k in LABEL_DISPLAY if k in buf and buf[k]]
    if active_keys:
        stat_cols = st.columns(len(active_keys))
        for i, key in enumerate(active_keys):
            with stat_cols[i]:
                st.markdown(
                    f'<div class="stat-box"><div class="num">{len(buf[key])}</div>'
                    f'<div class="lbl">{LABEL_DISPLAY[key]}</div></div>',
                    unsafe_allow_html=True,
                )
    st.markdown("")
    st.caption("✏️ 셀 클릭 → 직접 수정 / 맨 아래 빈 행 클릭 → 행 추가 / 행 선택 후 휴지통 아이콘 → 삭제")
    if "tc_parse_errors" in st.session_state and st.session_state["tc_parse_errors"]:
        with st.expander(f"⚠️ 일부 화면 TC 파싱 실패 ({len(st.session_state['tc_parse_errors'])}개) — 클릭해서 확인", expanded=True):
            for err in st.session_state["tc_parse_errors"]:
                st.markdown(f"**화면:** `{err['screen']}`")
                st.markdown(f"**오류:** `{err['error']}`")
                if err["raw"]:
                    st.code(err["raw"], language="text")
                st.markdown("---")
            st.caption("해당 화면은 TC가 생성되지 않았습니다. 다시 분석하거나 수동으로 추가해주세요.")
    any_dirty = False
    # ── UX / Policy / Edge 공통 렌더 ──
    for sec_key, meta in SECTION_META.items():
        if sec_key not in buf or not buf[sec_key]:
            continue
        st.markdown(f'<div class="section-title">{meta["label"]}</div>', unsafe_allow_html=True)
        df_orig = items_to_df(buf[sec_key], meta["cols"], meta["keys"])
        edited_df = st.data_editor(
            df_orig,
            use_container_width=True,
            hide_index=True,
            num_rows="dynamic",
            key=f"de_{sec_key}",
            # ▼ [변경 2] 높이: 기존 행 + 빈 행 1개만 표시
            height=calc_height(len(df_orig)),
            column_config=meta["col_config"],
        )
        if is_dirty(buf[sec_key], edited_df, meta["cols"], meta["keys"]):
            buf[sec_key] = df_to_items(edited_df, meta["cols"], meta["keys"])
            any_dirty = True
    # ── QA Test Cases ──
    if "test_cases" in buf and buf["test_cases"]:
        st.markdown('<div class="section-title">🧪 4. QA Test Cases</div>', unsafe_allow_html=True)
        tcs = buf["test_cases"]
        high_tcs = [t for t in tcs if t.get("priority") == "High"]
        med_tcs  = [t for t in tcs if t.get("priority") == "Medium"]
        low_tcs  = [t for t in tcs if t.get("priority") == "Low"]

        # ▼ Smoke/Regression은 Screen 열 제거, Full Coverage만 표시
        tc_strategy = st.session_state.get("current_tc_strategy", "regression")
        if tc_strategy == "full":
            _tc_cols = TC_COLS
            _tc_keys = TC_KEYS
            _tc_col_config = TC_COL_CONFIG
        else:
            _tc_cols = [c for c in TC_COLS if c != "Screen"]
            _tc_keys = [k for k in TC_KEYS if k != "screen"]
            _tc_col_config = {k: v for k, v in TC_COL_CONFIG.items() if k != "Screen"}

        q1, q2, q3, q4 = st.columns(4)
        with q1:
            st.markdown(f'<div class="stat-box"><div class="num">{len(tcs)}</div><div class="lbl">전체</div></div>', unsafe_allow_html=True)
        with q2:
            st.markdown(f'<div class="stat-box"><div class="num" style="color:#ef4444">{len(high_tcs)}</div><div class="lbl">🔴 High</div></div>', unsafe_allow_html=True)
        with q3:
            st.markdown(f'<div class="stat-box"><div class="num" style="color:#f59e0b">{len(med_tcs)}</div><div class="lbl">🟡 Medium</div></div>', unsafe_allow_html=True)
        with q4:
            st.markdown(f'<div class="stat-box"><div class="num" style="color:#10b981">{len(low_tcs)}</div><div class="lbl">🟢 Low</div></div>', unsafe_allow_html=True)
        st.markdown("")
        tab_all, tab_high, tab_med, tab_low = st.tabs([
            f"전체 ({len(tcs)})",
            f"🔴 High ({len(high_tcs)})",
            f"🟡 Medium ({len(med_tcs)})",
            f"🟢 Low ({len(low_tcs)})",
        ])
        def make_tc_df(tc_list):
            return items_to_df(tc_list, _tc_cols, _tc_keys)
        with tab_all:
            st.caption("✏️ 이 탭에서만 편집이 저장됩니다. 우선순위 탭은 읽기 전용입니다.")
            df_all = make_tc_df(tcs)
            edited_tc = st.data_editor(
                df_all,
                use_container_width=True,
                hide_index=True,
                num_rows="dynamic",
                key="de_tc",
                # ▼ [변경 2] 높이: 기존 행 + 빈 행 1개만 표시
                height=calc_height(len(df_all), max_height=1200),
                column_config=_tc_col_config,
            )
            if is_dirty(buf["test_cases"], edited_tc, _tc_cols, _tc_keys):
                buf["test_cases"] = df_to_items(edited_tc, _tc_cols, _tc_keys)
                any_dirty = True
        with tab_high:
            if high_tcs:
                st.dataframe(
                    make_tc_df(high_tcs),
                    use_container_width=True,
                    hide_index=True,
                    height=calc_height(len(high_tcs), max_height=600),
                    column_config=_tc_col_config,
                )
            else:
                st.info("High 항목이 없습니다.")
        with tab_med:
            if med_tcs:
                st.dataframe(
                    make_tc_df(med_tcs),
                    use_container_width=True,
                    hide_index=True,
                    height=calc_height(len(med_tcs), max_height=600),
                    column_config=_tc_col_config,
                )
            else:
                st.info("Medium 항목이 없습니다.")
        with tab_low:
            if low_tcs:
                st.dataframe(
                    make_tc_df(low_tcs),
                    use_container_width=True,
                    hide_index=True,
                    height=calc_height(len(low_tcs), max_height=600),
                    column_config=_tc_col_config,
                )
            else:
                st.info("Low 항목이 없습니다.")
    # ── 저장 & 다운로드 ──
    st.markdown("---")
    save_label = "💾 변경사항 저장 ●" if any_dirty else "💾 저장 (변경 없음)"
    save_type  = "primary" if any_dirty else "secondary"
    s1, s2, s3 = st.columns(3)
    with s1:
        if st.button(save_label, use_container_width=True, type=save_type, key="save_btn"):
            st.session_state.current_results = buf
            st.session_state[buf_key] = json.loads(json.dumps(buf))
            st.toast("편집 내용이 저장되었습니다 ✅")
    with s2:
        xlsx_data = build_xlsx_report(buf)
        st.download_button(
            "⬇️ Excel 다운로드 (.xlsx)",
            data=xlsx_data,
            file_name="ai_review.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
            key="dl_xlsx",
        )
    with s3:
        html_report = build_html_report(buf)
        st.download_button(
            "🌐 HTML 리포트 다운로드",
            data=html_report,
            file_name="ai_review.html",
            mime="text/html",
            use_container_width=True,
            key="dl_html",
        )
# ─────────────────────────────────────────────
# 분석 결과 표시
# ─────────────────────────────────────────────
if st.session_state.current_results:
    st.markdown("---")
    st.markdown('<div class="section-title">📊 분석 결과</div>', unsafe_allow_html=True)
    render_results(
        results=st.session_state.current_results,
        file_name=st.session_state.current_file_name,
    )
